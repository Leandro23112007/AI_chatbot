from flask import Flask, render_template, request, jsonify, send_from_directory, send_file
from flask_cors import CORS
import ollama
import uuid
import json
import os
import re
import bleach
import atexit
import threading
import logging
import time
import tempfile
import shutil
from duckduckgo_search import DDGS
from werkzeug.utils import secure_filename
from PIL import Image
from diffusers import StableDiffusionPipeline, StableDiffusionXLPipeline
import torch

app = Flask(__name__)
CORS(app)

CHATS_FILE = 'chats_data.json'
USER_INFO_FILE = 'user_infos.json'
SETTINGS_FILE = 'settings.json'
FEEDBACK_FILE = 'feedback_data.json'
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'txt'}

cancel_flags = {}

# Locks para garantir thread safety
chats_lock = threading.Lock()
user_infos_lock = threading.Lock()
settings_lock = threading.Lock()

SAVE_INTERVAL = 10  # segundos

# Configuração básica do logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[
        logging.FileHandler('app.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)

# Prompts comuns
PROMPT_EXTRAIR_INFO = (
    "Extraia e salve informações pessoais relevantes do usuário (nome, profissão, interesses, cidade, clube, etc) "
    "a partir da mensagem abaixo. Responda apenas com um JSON contendo os campos detectados.\n"
    "Mensagem: \"{user_text}\"\n"
    "Exemplo de resposta: {{\"nome\": \"Leandro\", \"clube\": \"Sporting\"}}"
)
PROMPT_EXTRAIR_MEMORIA = (
    "A partir da mensagem do usuário abaixo, extraia apenas uma frase curta e objetiva que resuma uma informação pessoal relevante sobre o usuário, "
    "como gostos, hobbies, interesses, línguas que fala, preferências, profissão, localidade, idade, etc. "
    "Ignore informações sobre tópicos de conversa, eventos históricos, notícias, ou qualquer coisa que não seja sobre o próprio usuário. "
    "Responda no formato: 'O utilizador ...'. Se não houver informação pessoal relevante, responda apenas com 'N/A'.\n"
    "Mensagem: \"{user_text}\""
)

def load_data(file_path, default_data=None):
    if default_data is None:
        default_data = {}
    if not os.path.exists(file_path):
        return default_data
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, IOError) as e:
        logging.warning(f"Falha ao carregar {file_path}: {e}")
        return default_data

def sanitize_filename(filename):
    # Permite apenas letras, números, hífen e sublinhado
    return re.sub(r'[^a-zA-Z0-9_.-]', '_', filename)

def save_data(file_path, data):
    # Caminho absoluto seguro
    file_path = os.path.abspath(file_path)
    dir_name = os.path.dirname(file_path)
    # Sanitiza o nome do arquivo (apenas para arquivos novos criados a partir de input)
    base = os.path.basename(file_path)
    safe_base = sanitize_filename(base)
    file_path = os.path.join(dir_name, safe_base)
    if not os.path.exists(dir_name):
        os.makedirs(dir_name, exist_ok=True)
    try:
        with tempfile.NamedTemporaryFile('w', encoding='utf-8', dir=dir_name, delete=False) as tf:
            json.dump(data, tf, ensure_ascii=False, indent=2)
            tf.flush()
            os.fsync(tf.fileno())
            tempname = tf.name
        os.replace(tempname, file_path)
        os.chmod(file_path, 0o600)  # Permissões restritas
    except Exception as e:
        logging.error(f"Erro ao salvar {file_path} (escrita segura): {e}")

class Chat:
    def __init__(self, cid, data):
        self.cid = cid
        self.data = data

    @property
    def name(self):
        return self.data.get('name', 'Novo Chat')

    @name.setter
    def name(self, value):
        self.data['name'] = value

    def add_message(self, user_text, ai_text):
        self.data.setdefault('conversation_blocks', []).append({
            'user_variants': [user_text],
            'ai_responses': [ai_text],
            'selected': 0
        })

    def get_conversation_blocks(self):
        return self.data.get('conversation_blocks', [])

    def to_dict(self):
        return self.data

# Helper to wrap chats dict with Chat objects
class ChatManager:
    def __init__(self, chats_dict):
        self.chats_dict = chats_dict
    def __getitem__(self, cid):
        return Chat(cid, self.chats_dict[cid])
    def __setitem__(self, cid, chat_obj):
        self.chats_dict[cid] = chat_obj.to_dict()
    def __delitem__(self, cid):
        del self.chats_dict[cid]
    def items(self):
        return ((cid, Chat(cid, data)) for cid, data in self.chats_dict.items())
    def get(self, cid):
        if cid in self.chats_dict:
            return Chat(cid, self.chats_dict[cid])
        return None
    def __contains__(self, cid):
        return cid in self.chats_dict

chats = ChatManager(load_data(CHATS_FILE))
user_infos = load_data(USER_INFO_FILE)

def get_default_settings():
    return {
        "theme": "dark",
        "language": "auto",
        "memory": {
            "reference_saved_memories": True,
            "reference_chat_history": False
        }
    }

settings = load_data(SETTINGS_FILE, get_default_settings())
if not os.path.exists(SETTINGS_FILE):
    save_data(SETTINGS_FILE, settings)

def new_chat_obj(name="Novo Chat"):
    return {'name': name, 'conversation_blocks': []}

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/chats', methods=['GET'])
def list_chats():
    return jsonify([{'id': cid, 'name': chat.name} for cid, chat in chats.items()])

@app.route('/chat', methods=['POST'])
def create_chat():
    cid = str(uuid.uuid4())
    with chats_lock:
        chats[cid] = Chat(cid, new_chat_obj())
    return jsonify({'id': cid, 'name': chats[cid].name})

@app.route('/chat/<cid>', methods=['GET'])
def get_chat(cid):
    chat = chats.get(cid)
    if not chat:
        return jsonify({'error': 'Chat não encontrado'}), 404
    return jsonify(chat.get_conversation_blocks())

@app.route('/chat/<cid>', methods=['DELETE'])
def delete_chat(cid):
    with chats_lock:
        if cid in chats:
            del chats[cid]
            return jsonify({'ok': True})
    return jsonify({'error': 'Chat não encontrado'}), 404

@app.route('/chat/<cid>/rename', methods=['POST'])
def rename_chat(cid):
    new_name = request.json.get('name', '').strip()
    new_name = bleach.clean(new_name)
    if len(new_name) > 200:
        return jsonify({'error': 'Nome demasiado longo (máx 200 caracteres)'}), 400
    with chats_lock:
        if cid in chats and new_name:
            chats[cid].name = new_name
            return jsonify({'ok': True})
    return jsonify({'error': 'Chat não encontrado ou nome inválido'}), 400

# Função genérica para executar prompts no Ollama, com retries exponenciais
def executar_prompt(prompt, model='llama3.1:8b', max_retries=3):
    delay = 1
    for attempt in range(max_retries):
        try:
            response = ollama.chat(model=model, messages=[{'role': 'system', 'content': prompt}])
            return response['message']['content']
        except Exception as e:
            logging.warning(f"Erro ao executar prompt no Ollama (tentativa {attempt+1}): {e}")
            if attempt < max_retries - 1:
                time.sleep(delay)
                delay *= 2
    return None

def extrair_info_com_ia(user_text, user_info):
    prompt = PROMPT_EXTRAIR_INFO.format(user_text=user_text)
    result = executar_prompt(prompt)
    if result:
        try:
            # Tentar encontrar o primeiro JSON válido na resposta
            match = re.search(r'\{.*\}', result, re.DOTALL)
            if match:
                extracted = json.loads(match.group(0))
                if isinstance(extracted, dict):
                    user_info.update(extracted)
                else:
                    logging.warning(f"JSON extraído não é dict: {extracted}")
            else:
                logging.warning(f"Nenhum JSON válido encontrado na resposta: {result}")
        except Exception as e:
            logging.warning(f"Erro ao interpretar JSON extraído: {e}\nConteúdo: {result}")
    return user_info

def extrair_memoria_resumida(user_text):
    cumprimentos = [
        'bom dia', 'boa tarde', 'boa noite', 'olá', 'ola', 'oi', 'tudo bem', 'alô', 'hello', 'hi',
        'boa madrugada', 'saudações', 'salve', 'e aí', 'eai', 'yo', 'hey', 'oiê', 'oii', 'oiii', 'al'
    ]
    texto = user_text.strip().lower()
    cumprimentos_regex = [
        r'^[a\slo]+$',   # alo, aloo, aaalllloooo
        r'^[o\si]+$',    # oi, oii, oiii
        r'^[o\sla]+$',   # ola, ollla, ooolaaa
        r'^[h\si]+$',    # hi, hii, hiii
        r'^[e\sai]+$',   # eai, eaiii
    ]
    for pattern in cumprimentos_regex:
        if re.fullmatch(pattern, texto):
            return None
    if texto in cumprimentos or len(texto) < 10:
        return None
    prompt = PROMPT_EXTRAIR_MEMORIA.format(user_text=user_text)
    frase = executar_prompt(prompt)
    if frase and frase.strip() != 'N/A':
        return frase.strip()
    return None

def gerar_nome_conversa_primeira_mensagem(msg):
    prompt = (
        "Gere um nome curto, técnico e descritivo para uma conversa de chat com base na mensagem abaixo. "
        "O nome deve ter no máximo 5 palavras, evitar adjetivos vagos ou genéricos e descrever claramente o tema da conversa. "
        "Se a mensagem for apenas uma saudação ou não indicar um tema, sugira um título neutro como 'Saudações' ou 'Cumprimentos'. "
        "Responda apenas com o nome sugerido, sem explicações ou pontuação extra.\n"
        f"Mensagem: \"{msg}\""
    )
    try:
        response = ollama.chat(model='llama3.1:8b', messages=[{'role': 'system', 'content': prompt}])
        nome = response['message']['content'].strip()
        return nome
    except Exception as e:
        logging.warning(f"Falha ao gerar nome da conversa: {e}")
        return msg[:30]  # fallback: primeiros 30 caracteres

@app.route('/chat/<cid>/cancel', methods=['POST'])
def cancel_chat_response(cid):
    cancel_flags[cid] = True
    return jsonify({'ok': True})

# Carregar o pipeline uma vez (global, para performance)
sd_pipe = None
def get_sd_pipe():
    global sd_pipe
    if sd_pipe is None:
        print("[DEBUG] Carregando modelo SDXL para CPU...")
        sd_pipe = StableDiffusionXLPipeline.from_pretrained(
            "stabilityai/stable-diffusion-xl-base-1.0",
            torch_dtype=torch.float32
        )
        sd_pipe = sd_pipe.to("cpu")
    return sd_pipe

def is_image_request(user_text):
    # Detecta pedidos de imagem de forma simples
    keywords = [
        "gera uma imagem", "cria uma imagem", "desenha", "faz uma imagem", "ilustra", "imagem de", "picture of", "draw", "generate an image", "create an image"
    ]
    user_text_lower = user_text.lower()
    return any(kw in user_text_lower for kw in keywords)

def generate_image_from_text(prompt, output_path):
    pipe = get_sd_pipe()
    print(f"[DEBUG] Gerando imagem para prompt: {prompt}")
    image = pipe(prompt).images[0]
    image.save(output_path)
    print(f"[DEBUG] Imagem salva em: {output_path}")

@app.route('/chat/<cid>/send', methods=['POST'])
def send_message(cid):
    user_text = request.json.get('message')
    if not user_text:
        return jsonify({'error': 'Mensagem vazia'}), 400
    user_text = bleach.clean(user_text)
    if len(user_text) > 20000:
        return jsonify({'error': 'Mensagem demasiado longa (máx 10000 caracteres)'}), 400
    if cid not in chats:
        return jsonify({'error': 'Chat não encontrado'}), 404

    # Salvar imediatamente o bloco do usuário com resposta AI como None
    with chats_lock:
        chats[cid].data.setdefault('conversation_blocks', []).append({
            'user_variants': [user_text],
            'ai_responses': [None],
            'selected': 0
        })
        save_data(CHATS_FILE, chats.chats_dict)

    # 1. Construir o contexto do sistema
    system_messages = []
    system_message = """
Identidade e Função:
És o Llama 3, um modelo de linguagem do Ollama. Foste criado para ajudar os utilizadores com informações úteis, gerar texto, responder a perguntas, resolver problemas e muito mais.

Estilo de Comunicação:
Sê direto, claro e profissional.
Não faças elogios exagerados ou bajulação.
Mantém um tom humano e caloroso, mas evita soar robótico ou artificialmente entusiasmado.
Prioriza a honestidade: se não souberes algo, diz-o abertamente.
Evita adivinhar o que o utilizador quer sem confirmação. Pergunta se necessário.

Linguagem e Preferências do Utilizador:
Lê e adapta-te às preferências do utilizador com base no contexto da conversa (por exemplo, se preferir respostas curtas, em português europeu, técnicas, etc.).

Privacidade e Segurança:
Nunca peças, guardes ou reveles dados sensíveis.
"""
    system_messages.append({'role': 'system', 'content': system_message})
    user_id = "default_user"
    current_user_info = user_infos.get(user_id, {})
    updated_info = extrair_info_com_ia(user_text, dict(current_user_info))
    if updated_info != current_user_info:
        with user_infos_lock:
            user_infos[user_id] = updated_info
    memoria_foi_atualizada = False
    frase_memoria = extrair_memoria_resumida(user_text)
    if frase_memoria:
        with user_infos_lock:
            user_infos.setdefault(user_id, {}).setdefault('memorias_resumidas', [])
            if frase_memoria not in user_infos[user_id]['memorias_resumidas']:
                user_infos[user_id]['memorias_resumidas'].append(frase_memoria)
                # Limitar para as 20 mais recentes
                user_infos[user_id]['memorias_resumidas'] = user_infos[user_id]['memorias_resumidas'][-20:]
                memoria_foi_atualizada = True
    # Adicionar info do usuário se a configuração permitir
    if settings.get('memory', {}).get('reference_saved_memories', False) and updated_info:
        system_context = "As informações a seguir são sobre o usuário. Use-as para personalizar suas respostas de forma implícita, sem mencioná-las diretamente.\n"
        for key, value in updated_info.items():
            if value:
                system_context += f"- {key.replace('_', ' ').capitalize()}: {value}\n"
        # Adicionar frases resumidas
        memorias_resumidas = user_infos.get(user_id, {}).get('memorias_resumidas', [])
        if memorias_resumidas:
            system_context += "- Outras informações relevantes:\n"
            for frase in memorias_resumidas:
                system_context += f"  - {frase}\n"
        system_messages.append({'role': 'system', 'content': system_context})

    # 2. Construir o histórico da conversa
    history_messages = []
    if settings.get('memory', {}).get('reference_chat_history', False):
        for block in chats[cid].get_conversation_blocks():
            history_messages.append({'role': 'user', 'content': block['user_variants'][block['selected']]})
            history_messages.append({'role': 'assistant', 'content': block['ai_responses'][block['selected']]})

    # 3. Montar a mensagem final para o Ollama
    ollama_messages = system_messages + history_messages + [{'role': 'user', 'content': user_text}]

    # --- NOVO FLUXO: verificação de geração de ficheiro ---
    generated_filename = None
    file_type = None
    
    # Primeiro, gera a resposta da IA normalmente
    try:
        response = ollama.chat(model='llama3.1:8b', messages=ollama_messages)
        ai_text = response['message']['content']
    except Exception as e:
        logging.error(f"Erro ao conectar ao Ollama: {e}")
        ai_text = "Lamento, não consigo processar isso agora."

    # Depois, verifica se o usuário pediu um ficheiro e gera o ficheiro se necessário
    if user_wants_file(user_text):
        file_type = get_file_type(user_text)
        if not file_type or file_type not in ['pdf', 'txt', 'csv', 'json', 'py', 'docx', 'xlsx', 'pptx']:
            ai_text = "Por favor, indique o tipo de ficheiro (ex: pdf, docx, txt, etc.)"
            with chats_lock:
                chats[cid].data['conversation_blocks'][-1]['ai_responses'][0] = ai_text
                save_data(CHATS_FILE, chats.chats_dict)
            return jsonify({
                'ai_text': ai_text,
                'block_idx': len(chats[cid].get_conversation_blocks()) - 1,
                'memoria_atualizada': False
            })
        generated_filename = ai_generate_file_if_requested(user_text, ai_text, cid, file_type)
        if generated_filename:
            download_url = f"/files/{sanitize_filename(cid)}/{generated_filename}"
            ai_text += f"\n\n[Download do arquivo gerado]({download_url})"
            # Atualizar o bloco com a resposta do AI
            with chats_lock:
                chats[cid].data['conversation_blocks'][-1]['ai_responses'][0] = ai_text
                save_data(CHATS_FILE, chats.chats_dict)
            return jsonify({
                'ai_text': ai_text,
                'block_idx': len(chats[cid].get_conversation_blocks()) - 1,
                'memoria_atualizada': True
            })

    # --- SUPORTE À GERAÇÃO DE IMAGEM ---
    if is_image_request(user_text):
        chat_folder = os.path.join(UPLOAD_FOLDER, sanitize_filename(cid))
        os.makedirs(chat_folder, exist_ok=True)
        filename = f"imagem_{int(time.time())}.png"
        file_path = os.path.join(chat_folder, filename)
        generate_image_from_text(user_text, file_path)
        download_url = f"/files/{sanitize_filename(cid)}/{filename}"
        ai_text = f"Imagem gerada com sucesso! [Download da imagem gerada]({download_url})"
        with chats_lock:
            chats[cid].data['conversation_blocks'][-1]['ai_responses'][0] = ai_text
            save_data(CHATS_FILE, chats.chats_dict)
        return jsonify({
            'ai_text': ai_text,
            'block_idx': len(chats[cid].get_conversation_blocks()) - 1,
            'memoria_atualizada': False
        })

    # --- CANCELAMENTO: se o chat foi cancelado, não salvar resposta ---
    if cancel_flags.get(cid):
        cancel_flags.pop(cid, None)
        # Atualiza o bloco com resposta de cancelamento
        with chats_lock:
            chats[cid].data['conversation_blocks'][-1]['ai_responses'][0] = '⏹️ Resposta cancelada pelo usuário.'
            save_data(CHATS_FILE, chats.chats_dict)
        return jsonify({'ai_text': '⏹️ Resposta cancelada pelo usuário.', 'memoria_atualizada': memoria_foi_atualizada})

    # Atualizar o bloco com a resposta do AI
    with chats_lock:
        chats[cid].data['conversation_blocks'][-1]['ai_responses'][0] = ai_text
        # NOVO: se for a primeira mensagem, gerar nome automático
        if len(chats[cid].get_conversation_blocks()) == 1:
            nome_sugerido = gerar_nome_conversa_primeira_mensagem(user_text)
            if nome_sugerido:
                chats[cid].name = nome_sugerido
        save_data(CHATS_FILE, chats.chats_dict)
    return jsonify({
        'ai_text': ai_text, 
        'block_idx': len(chats[cid].get_conversation_blocks()) - 1,
        'memoria_atualizada': memoria_foi_atualizada
    })

@app.route('/user_info', methods=['GET'])
def get_user_info_endpoint():
    user_id = "default_user"
    return jsonify(user_infos.get(user_id, {}))

@app.route('/user_info', methods=['POST'])
def update_user_info_endpoint():
    user_id = "default_user"
    sanitized_info = {}
    for k, v in request.json.items():
        if isinstance(v, str):
            sanitized_info[k] = bleach.clean(v)[:200]
        elif isinstance(v, list):
            sanitized_info[k] = [bleach.clean(str(item))[:200] for item in v]
        else:
            sanitized_info[k] = v
    with user_infos_lock:
        user_infos[user_id] = sanitized_info
    return jsonify({'ok': True})

@app.route('/api/settings', methods=['GET'])
def get_settings_endpoint():
    global settings
    settings = load_data(SETTINGS_FILE, get_default_settings())
    return jsonify(settings)

@app.route('/api/settings', methods=['POST'])
def update_settings_endpoint():
    global settings
    with settings_lock:
        settings = request.json
    return jsonify({"status": "success"})

@app.route('/user_memory/delete', methods=['POST'])
def delete_user_memory():
    user_id = "default_user"
    idx = request.json.get('idx')
    if user_id in user_infos and 'memorias_resumidas' in user_infos[user_id]:
        try:
            idx = int(idx)
            if 0 <= idx < len(user_infos[user_id]['memorias_resumidas']):
                user_infos[user_id]['memorias_resumidas'].pop(idx)
                save_data(USER_INFO_FILE, user_infos)
                return jsonify({'ok': True})
        except Exception as e:
            logging.warning(f"Erro ao apagar memória do usuário: {e}")
    return jsonify({'error': 'Índice inválido'}), 400

@app.route('/user_memory/delete_all', methods=['POST'])
def delete_all_user_memories():
    user_id = "default_user"
    if user_id in user_infos and 'memorias_resumidas' in user_infos[user_id]:
        user_infos[user_id]['memorias_resumidas'] = []
        save_data(USER_INFO_FILE, user_infos)
        return jsonify({'ok': True})
    return jsonify({'error': 'Nada para apagar'}), 400

def periodic_save():
    with chats_lock:
        save_data(CHATS_FILE, chats.chats_dict)
    with user_infos_lock:
        save_data(USER_INFO_FILE, user_infos)
    with settings_lock:
        save_data(SETTINGS_FILE, settings)
    # Reagendar
    global save_timer
    save_timer = threading.Timer(SAVE_INTERVAL, periodic_save)
    save_timer.daemon = True
    save_timer.start()

def save_all():
    with chats_lock:
        save_data(CHATS_FILE, chats.chats_dict)
    with user_infos_lock:
        save_data(USER_INFO_FILE, user_infos)
    with settings_lock:
        save_data(SETTINGS_FILE, settings)

save_timer = threading.Timer(SAVE_INTERVAL, periodic_save)
save_timer.daemon = True
save_timer.start()
atexit.register(save_all)

def load_feedback():
    if not os.path.exists(FEEDBACK_FILE):
        return {}
    try:
        with open(FEEDBACK_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def save_feedback(data):
    with open(FEEDBACK_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

@app.route('/feedback', methods=['POST'])
def add_feedback():
    data = request.json
    user_id = data.get('user_id', 'default_user')
    tema = data.get('tema', 'geral')
    feedback = data.get('feedback')  # 'positivo' ou 'negativo'
    comentario = data.get('comentario', '')
    mensagem_usuario = data.get('mensagem_usuario', '')
    resposta_ai = data.get('resposta_ai', '')
    timestamp = time.strftime('%Y-%m-%d %H:%M:%S')

    all_feedback = load_feedback()
    tema_feedbacks = all_feedback.setdefault(user_id, {}).setdefault(tema, [])

    # Verifica se já existe feedback para esta mensagem/resposta
    found = False
    for item in tema_feedbacks:
        if (item["mensagem_usuario"] == mensagem_usuario and
            item["resposta_ai"] == resposta_ai):
            # Atualiza o feedback existente
            item["feedback"] = feedback
            item["comentario"] = comentario
            item["timestamp"] = timestamp
            found = True
            break

    if not found:
        tema_feedbacks.append({
            "mensagem_usuario": mensagem_usuario,
            "resposta_ai": resposta_ai,
            "feedback": feedback,
            "comentario": comentario,
            "timestamp": timestamp
        })

    save_feedback(all_feedback)
    return jsonify({"ok": True})

@app.route('/feedback/stats', methods=['GET'])
def feedback_stats():
    all_feedback = load_feedback()
    stats = {}
    for user, temas in all_feedback.items():
        for tema, items in temas.items():
            stats.setdefault(tema, {"positivo": 0, "negativo": 0})
            for item in items:
                if item["feedback"] == "positivo":
                    stats[tema]["positivo"] += 1
                elif item["feedback"] == "negativo":
                    stats[tema]["negativo"] += 1
    return jsonify(stats)

def buscar_na_web(query, max_results=3):
    with DDGS() as ddgs:
        results = ddgs.text(query)
        return [r['body'] for r in results][:max_results]

@app.route('/api/websearch', methods=['POST'])
def websearch():
    data = request.json
    query = data.get('query', '')
    if not query or len(query) < 3:
        return jsonify({'error': 'Query muito curta'}), 400
    results = buscar_na_web(query)
    return jsonify({'results': results})

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/api/upload', methods=['POST'])
def upload_file():
    chat_id = request.form.get('chat_id')
    if not chat_id:
        return jsonify({'error': 'chat_id é obrigatório'}), 400
    if 'file' not in request.files:
        return jsonify({'error': 'Nenhum arquivo enviado'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'Nome de arquivo vazio'}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        chat_folder = os.path.join(UPLOAD_FOLDER, sanitize_filename(chat_id))
        os.makedirs(chat_folder, exist_ok=True)
        file.save(os.path.join(chat_folder, filename))
        return jsonify({'ok': True, 'filename': filename})
    return jsonify({'error': 'Tipo de arquivo não permitido'}), 400

@app.route('/api/files', methods=['GET'])
def list_files():
    chat_id = request.args.get('chat_id')
    if not chat_id:
        return jsonify({'error': 'chat_id é obrigatório'}), 400
    chat_folder = os.path.join(UPLOAD_FOLDER, sanitize_filename(chat_id))
    if not os.path.exists(chat_folder):
        return jsonify({'files': []})
    files = [f for f in os.listdir(chat_folder) if allowed_file(f)]
    return jsonify({'files': files})

@app.route('/api/download', methods=['GET'])
def download_file():
    chat_id = request.args.get('chat_id')
    filename = request.args.get('filename')
    if not chat_id or not filename:
        return jsonify({'error': 'chat_id e filename são obrigatórios'}), 400
    chat_folder = os.path.join(UPLOAD_FOLDER, sanitize_filename(chat_id))
    if not os.path.exists(chat_folder):
        return jsonify({'error': 'Arquivo não encontrado'}), 404
    filename = secure_filename(filename)
    file_path = os.path.join(chat_folder, filename)
    if not os.path.exists(file_path):
        return jsonify({'error': 'Arquivo não encontrado'}), 404
    return send_from_directory(chat_folder, filename, as_attachment=True)

def write_markdown_line(pdf, line):
    # Divide a linha em partes: texto normal e **negrito**
    parts = re.split(r'(\*\*.*?\*\*)', line)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            pdf.set_font('Arial', 'B', 12)
            pdf.write(8, part[2:-2])  # Remove os **
            pdf.set_font('Arial', '', 12)
        else:
            pdf.write(8, part)
    pdf.ln(8)

# --- AI gera arquivo sob pedido ---
def ai_generate_file_if_requested(user_text, ai_text, chat_id, file_type=None):
    import io
    import fpdf
    import csv
    import json as pyjson
    import pandas as pd
    from docx import Document
    from pptx import Presentation
    import openpyxl
    lower = user_text.lower()
    chat_folder = os.path.join(UPLOAD_FOLDER, sanitize_filename(chat_id))
    os.makedirs(chat_folder, exist_ok=True)
    content = user_text.split('com', 1)[-1].strip() if 'com' in user_text else ai_text
    filename = None
    try:
        if file_type == 'pdf':
            print(f"[DEBUG] Iniciando criação de PDF para chat_id={chat_id} na pasta {chat_folder}")
            pdf = fpdf.FPDF()
            pdf.add_page()
            pdf.set_font('Arial', size=12)
            for line in (content if content else ai_text).split('\n'):
                write_markdown_line(pdf, line)
            filename = f'gerado_{int(time.time())}.pdf'
            file_path = os.path.join(chat_folder, filename)
            pdf.output(file_path)
            print(f"[DEBUG] PDF criado com sucesso: {file_path}")
            logging.info(f"PDF criado: {file_path}")
        elif file_type == 'txt':
            filename = f'gerado_{int(time.time())}.txt'
            file_path = os.path.join(chat_folder, filename)
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content if content else ai_text)
        elif file_type == 'csv':
            filename = f'gerado_{int(time.time())}.csv'
            file_path = os.path.join(chat_folder, filename)
            rows = [row.split(';') for row in (content if content else ai_text).split('\n') if row.strip()]
            with open(file_path, 'w', encoding='utf-8', newline='') as f:
                writer = csv.writer(f)
                writer.writerows(rows)
        elif file_type == 'json':
            filename = f'gerado_{int(time.time())}.json'
            file_path = os.path.join(chat_folder, filename)
            try:
                data = pyjson.loads(content)
            except Exception:
                data = {'conteudo': content if content else ai_text}
            with open(file_path, 'w', encoding='utf-8') as f:
                pyjson.dump(data, f, ensure_ascii=False, indent=2)
        elif file_type == 'py':
            filename = f'gerado_{int(time.time())}.py'
            file_path = os.path.join(chat_folder, filename)
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content if content else ai_text)
        elif file_type == 'docx':
            filename = f'gerado_{int(time.time())}.docx'
            file_path = os.path.join(chat_folder, filename)
            doc = Document()
            doc.add_paragraph(content if content else ai_text)
            doc.save(file_path)
        elif file_type == 'xlsx':
            filename = f'gerado_{int(time.time())}.xlsx'
            file_path = os.path.join(chat_folder, filename)
            try:
                df = pd.read_csv(io.StringIO(content))
            except Exception:
                df = pd.DataFrame({'Conteudo': [(content if content else ai_text)]})
            df.to_excel(file_path, index=False)
        elif file_type == 'pptx':
            filename = f'gerado_{int(time.time())}.pptx'
            file_path = os.path.join(chat_folder, filename)
            ppt = Presentation()
            slide = ppt.slides.add_slide(ppt.slide_layouts[0])
            slide.shapes.title.text = (content if content else ai_text)[:100]
            ppt.save(file_path)
        else:
            logging.warning(f"Tipo de ficheiro não suportado: {file_type}")
            return None
        logging.info(f"Ficheiro gerado: {file_path}")
        print(f"[DEBUG] Ficheiro gerado: {file_path}")
        return filename
    except Exception as e:
        logging.error(f"Erro ao gerar ficheiro {file_type}: {e}")
        print(f"[DEBUG] Erro ao gerar ficheiro {file_type}: {e}")
        return None

def get_file_type(user_text):
    user_text_lower = user_text.lower()
    for ext in ['pdf', 'txt', 'csv', 'json', 'py', 'docx', 'xlsx', 'pptx']:
        if ext in user_text_lower:
            print(f"[get_file_type] Detetado tipo '{ext}' diretamente no texto do utilizador.")
            return ext
    import time
    prompt = (
        "Que tipo de ficheiro é suposto ser gerado? Responde só pdf, txt, csv, json, py, docx, xlsx, pptx ou outro tipo simples, sem contexto nem formatação. "
        f"Frase: {user_text}"
    )
    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = ollama.chat(model='llama3.1:8b', messages=[{'role': 'system', 'content': prompt}])
            answer = response['message']['content'].strip().lower()
            print(f"[AI get_file_type] Tentativa {attempt+1}: Pergunta: {prompt}\nResposta: {answer}")
            for ext in ['pdf', 'txt', 'csv', 'json', 'py', 'docx', 'xlsx', 'pptx']:
                if ext in answer:
                    return ext
            if attempt == 1:
                # Após a segunda tentativa, se não encontrou, retorna vazio especial
                return ''
            if answer:
                return answer
            else:
                print(f"[AI get_file_type] Resposta vazia, retry {attempt+1}")
                time.sleep(1)
        except Exception as e:
            logging.warning(f"Erro ao obter tipo de ficheiro: {e}")
            print(f"[AI get_file_type] Erro na tentativa {attempt+1}: {e}")
            time.sleep(1)
    return ''

def user_wants_file(user_text):
    import time
    # Padrões explícitos de geração de ficheiro
    explicit_patterns = [
        r'^(gera|cria|exporta|envia|faz|quero baixar|quero um|quero uma|quero o|quero a).*\b(pdf|ficheiro|arquivo|documento|word|excel|pptx|csv|json|txt|py)\b',
    ]
    user_text_lower = user_text.lower()
    for pattern in explicit_patterns:
        if re.search(pattern, user_text_lower):
            print(f"[user_wants_file] Detetado padrão explícito de geração de ficheiro: {user_text}")
            return True

    prompt = (
        "Isto é para gerar um ficheiro? Responde só sim ou não, sem contexto nem formatação. "
        f"Frase: {user_text}"
    )
    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = ollama.chat(model='llama3.1:8b', messages=[{'role': 'system', 'content': prompt}])
            answer = response['message']['content'].strip().lower()
            print(f"[AI user_wants_file] Tentativa {attempt+1}: Pergunta: {prompt}\nResposta: {answer}")
            if answer:
                if (
                    answer.startswith('sim') or
                    'sim' in answer or
                    'pdf' in answer or
                    'ficheiro' in answer or
                    'arquivo' in answer or
                    'documento' in answer or
                    'word' in answer or
                    'excel' in answer or
                    'pptx' in answer or
                    'csv' in answer or
                    'json' in answer or
                    'txt' in answer or
                    'py' in answer
                ):
                    return True
                return False
            else:
                print(f"[AI user_wants_file] Resposta vazia, retry {attempt+1}")
                time.sleep(1)
        except Exception as e:
            logging.warning(f"Erro ao decidir se deve gerar ficheiro: {e}")
            print(f"[AI user_wants_file] Erro na tentativa {attempt+1}: {e}")
            time.sleep(1)
    return False

@app.route('/files/<chat_id>/<filename>')
def download_file_simple(chat_id, filename):
    chat_folder = os.path.join(UPLOAD_FOLDER, sanitize_filename(chat_id))
    filename = secure_filename(filename)
    return send_from_directory(chat_folder, filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True, port=5000)